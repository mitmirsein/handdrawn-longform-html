#!/usr/bin/env python3
"""Render a reviewed deck.json into a local 16:9 PPTX.

This renderer only assembles local text and images. It does not research,
generate images, or fill in missing claims.
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


W, H, M = 13.333, 7.5, 0.62
INK, BODY, ACCENT, GREY = (26, 26, 26), (51, 51, 51), (226, 74, 59), (150, 150, 150)


def _box(slide, x, y, w, h, Inches):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    return frame


def _write(frame, text, *, font, size, color, Pt, RGBColor, bold=False, align=None, line=1.2):
    # python-pptx writes the requested typeface as a:latin only. LibreOffice
    # otherwise falls back to a sans-serif face for Korean glyphs, so also set
    # the East Asian and complex-script slots explicitly.
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement

    for index, value in enumerate(str(text).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        if align is not None:
            paragraph.alignment = align
        paragraph.line_spacing = line
        run = paragraph.add_run()
        run.text = value
        run.font.name = font
        run.font.size = Pt(size)
        # The installed Korean brush/pen faces are regular-only. Asking for a
        # synthetic bold face makes LibreOffice fall back to a modern sans face,
        # which defeats the rough diary treatment.
        run.font.bold = False if font in {"Nanum Pen Script", "Nanum Brush Script", "Gaegu", "BM Hanna Pro", "BM Hanna Air", "BM Yeonsung"} else bold
        run.font.color.rgb = RGBColor(*color)
        rpr = run._r.get_or_add_rPr()
        for tag in ("a:ea", "a:cs"):
            element = rpr.find(qn(tag))
            if element is None:
                element = OxmlElement(tag)
                rpr.append(element)
            element.set("typeface", font)


def _ratio(path: Path) -> float:
    try:
        from PIL import Image
        with Image.open(path) as image:
            return image.width / image.height
    except (ImportError, OSError, ValueError, ZeroDivisionError):
        return 16 / 9


def _prepared(path: Path, max_width: int, cache: Path) -> Path:
    if max_width <= 0:
        return path
    try:
        from PIL import Image
    except ImportError:
        return path
    try:
        with Image.open(path) as image:
            if image.width <= max_width:
                return path
            output = cache / f"{path.stem}__{max_width}.png"
            if output.is_file() and output.stat().st_mtime_ns >= path.stat().st_mtime_ns:
                return output
            cache.mkdir(parents=True, exist_ok=True)
            resampling = getattr(Image, "Resampling", Image)
            resized = image.resize((max_width, max(1, int(image.height * max_width / image.width))), resampling.LANCZOS)
            resized.save(output, "PNG", optimize=True)
            return output
    except (OSError, ValueError):
        return path


def _image(slide, path: Path, x, y, w, h, Inches):
    ratio = _ratio(path)
    if w / h > ratio:
        ih, iw = h, h * ratio
    else:
        iw, ih = w, w / ratio
    slide.shapes.add_picture(str(path), Inches(x + (w - iw) / 2), Inches(y + (h - ih) / 2), width=Inches(iw), height=Inches(ih))


def _source(slide, item, body_font, Inches, Pt, RGBColor):
    source = item.get("source")
    lines = item.get("source_lines")
    if lines:
        suffix = f" · 원문 {min(lines)}–{max(lines)}행"
        source = f"{source or ''}{suffix}".strip(" ·")
    if source:
        frame = _box(slide, M, 6.78, W - 2 * M, 0.32, Inches)
        _write(frame, source, font=body_font, size=8.5, color=GREY, Pt=Pt, RGBColor=RGBColor, line=1.05)


def _notes(slide, item):
    notes = item.get("speaker_notes") or item.get("notes")
    if notes:
        slide.notes_slide.notes_text_frame.text = str(notes)


def _card(slide, text, x, y, w, h, fonts, *, accent=False):
    rough = fonts.get("rough", False)
    shape_type = fonts["MSO_SHAPE"].RECTANGLE if rough else fonts["MSO_SHAPE"].ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, fonts["Inches"](x), fonts["Inches"](y), fonts["Inches"](w), fonts["Inches"](h))
    shape.fill.solid()
    if rough:
        # Keep the reference's paper-like white space and thin pen outlines.
        shape.fill.fore_color.rgb = fonts["RGBColor"](*(255, 248, 245) if accent else (255, 255, 255))
        shape.line.color.rgb = fonts["RGBColor"](*(226, 74, 59) if accent else (26, 26, 26))
        shape.line.width = fonts["Pt"](1.05)
    else:
        shape.fill.fore_color.rgb = fonts["RGBColor"](*(255, 244, 241) if accent else (245, 247, 249))
        shape.line.color.rgb = fonts["RGBColor"](*(226, 74, 59) if accent else (190, 196, 202))
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = fonts["Inches"](0.14)
    frame.margin_top = frame.margin_bottom = fonts["Inches"](0.08)
    _write(frame, text, font=fonts["font"], size=20 if rough else 18, color=INK, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], bold=accent, align=fonts["PP_ALIGN"].CENTER)


def _labels(value, separators):
    text = str(value or "").replace("↔", "|")
    for separator in separators:
        text = text.replace(separator, "|")
    return [part.strip() for part in text.split("|") if part.strip()]


def _headline(slide, item, font, Pt, RGBColor, Inches, *, size=35, y=0.92, align=None):
    frame = _box(slide, M, y, W - 2 * M, 1.1, Inches)
    _write(frame, item.get("headline", ""), font=font, size=size, color=INK, Pt=Pt, RGBColor=RGBColor, bold=True, align=align, line=1.05)


def _layout_title(slide, item, image, fonts):
    frame = _box(slide, M, 1.9, 6.05, 2.7, fonts["Inches"])
    _write(frame, item.get("headline", ""), font=fonts["font"], size=46, color=INK, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], bold=True, line=1.08)
    if item.get("body") or item.get("visible_text"):
        body = _box(slide, M, 5.0, 6.05, 0.9, fonts["Inches"])
        _write(body, item.get("body") or item.get("visible_text"), font=fonts["font"], size=18, color=BODY, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], line=1.35)
    if image:
        _image(slide, image, 6.9, 1.25, 5.8, 4.95, fonts["Inches"])


def _layout_content(slide, item, image, fonts):
    _headline(slide, item, fonts["font"], fonts["Pt"], fonts["RGBColor"], fonts["Inches"])
    if image:
        _image(slide, image, 5.5, 1.85, 7.2, 4.7, fonts["Inches"])
        width = 4.35
    else:
        width = W - 2 * M
    text = item.get("body") or item.get("visible_text") or item.get("one_sentence_takeaway", "")
    if text:
        frame = _box(slide, M, 2.0, width, 3.8, fonts["Inches"])
        _write(frame, text, font=fonts["font"], size=18, color=BODY, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], line=1.4)


def _layout_full(slide, item, image, fonts):
    _headline(slide, item, fonts["font"], fonts["Pt"], fonts["RGBColor"], fonts["Inches"], y=0.8)
    if image:
        _image(slide, image, M, 1.9, W - 2 * M, 4.25, fonts["Inches"])
    text = item.get("body") or item.get("visible_text") or item.get("one_sentence_takeaway", "")
    if text:
        frame = _box(slide, M, 5.95, W - 2 * M, 0.55, fonts["Inches"])
        _write(frame, text, font=fonts["font"], size=17, color=BODY, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], align=fonts["PP_ALIGN"].CENTER)


def _layout_quote(slide, item, image, fonts):
    frame = _box(slide, 1.45, 2.25, 10.45, 2.4, fonts["Inches"])
    _write(frame, item.get("headline", ""), font=fonts["font"], size=42, color=INK, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], bold=True, align=fonts["PP_ALIGN"].CENTER, line=1.14)
    text = item.get("body") or item.get("visible_text")
    if text:
        body_width = 7.55 if image else 10.45
        body = _box(slide, 1.45, 4.8, body_width, 0.8, fonts["Inches"])
        _write(body, text, font=fonts["font"], size=18, color=BODY, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], align=fonts["PP_ALIGN"].CENTER)
    if image:
        _image(slide, image, 9.35, 4.55, 2.9, 1.85, fonts["Inches"])


def _layout_comparison(slide, item, image, fonts):
    _headline(slide, item, fonts["font"], fonts["Pt"], fonts["RGBColor"], fonts["Inches"])
    labels = _labels(item.get("visible_text") or item.get("one_sentence_takeaway"), ["+", "↔", "vs", " VS "])
    labels = (labels + ["비교 기준"])[:2] or ["비교 대상", "비교 기준"]
    _card(slide, labels[0], M, 2.15, 5.75, 2.15, fonts, accent=True)
    _card(slide, labels[1], 6.95, 2.15, 5.75, 2.15, fonts)
    body = item.get("body") or item.get("one_sentence_takeaway")
    if body:
        body_width = 8.2 if image else W - 2 * M
        frame = _box(slide, M, 4.75, body_width, 1.0, fonts["Inches"])
        _write(frame, body, font=fonts["font"], size=16, color=BODY, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], align=fonts["PP_ALIGN"].CENTER)
    if image:
        _image(slide, image, 9.35, 4.5, 2.9, 1.95, fonts["Inches"])


def _layout_timeline(slide, item, image, fonts):
    _headline(slide, item, fonts["font"], fonts["Pt"], fonts["RGBColor"], fonts["Inches"])
    labels = _labels(item.get("visible_text") or item.get("one_sentence_takeaway"), ["·", "→"])
    labels = labels[:5] or ["배경", "전환", "결론"]
    width = (W - 2 * M - 0.25 * (len(labels) - 1)) / len(labels)
    for index, label in enumerate(labels):
        _card(slide, label, M + index * (width + 0.25), 2.55, width, 1.35, fonts, accent=index == len(labels) - 1)
        if index < len(labels) - 1:
            arrow = _box(slide, M + (index + 1) * width + index * 0.25 - 0.01, 3.0, 0.27, 0.35, fonts["Inches"])
            _write(arrow, "→", font=fonts["font"], size=22, color=ACCENT, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], align=fonts["PP_ALIGN"].CENTER)
    if image:
        _image(slide, image, 9.35, 4.45, 2.9, 1.95, fonts["Inches"])


def _layout_flow(slide, item, image, fonts):
    _headline(slide, item, fonts["font"], fonts["Pt"], fonts["RGBColor"], fonts["Inches"])
    labels = _labels(item.get("visible_text") or item.get("one_sentence_takeaway"), ["→"])
    labels = labels[:5] or ["원인", "변화", "결과"]
    width = min(2.25, (W - 2 * M - 0.3 * (len(labels) - 1)) / len(labels))
    left = (W - (len(labels) * width + (len(labels) - 1) * 0.3)) / 2
    for index, label in enumerate(labels):
        _card(slide, label, left + index * (width + 0.3), 2.45, width, 1.45, fonts, accent=index == len(labels) - 1)
        if index < len(labels) - 1:
            arrow = _box(slide, left + (index + 1) * width + index * 0.3 - 0.02, 2.95, 0.3, 0.35, fonts["Inches"])
            _write(arrow, "→", font=fonts["font"], size=22, color=ACCENT, Pt=fonts["Pt"], RGBColor=fonts["RGBColor"], align=fonts["PP_ALIGN"].CENTER)
    if image:
        _image(slide, image, 9.35, 4.45, 2.9, 1.95, fonts["Inches"])


def _layout_table(slide, item, image, fonts):
    _headline(slide, item, fonts["font"], fonts["Pt"], fonts["RGBColor"], fonts["Inches"])
    labels = _labels(item.get("visible_text") or item.get("one_sentence_takeaway"), ["|", ";", "·"])
    labels = labels[:6] or ["핵심", "근거", "적용"]
    for index, label in enumerate(labels):
        row = index // 3
        column = index % 3
        _card(slide, label, M + column * 4.1, 2.0 + row * 1.65, 3.75, 1.25, fonts, accent=index == 0)
    if image:
        _image(slide, image, 9.35, 5.05, 2.9, 1.35, fonts["Inches"])


def render(deck_path: Path, output: Path, max_width: int = 1800) -> Path:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError(f"python-pptx is required: {sys.executable} -m pip install python-pptx") from exc

    deck = json.loads(deck_path.read_text(encoding="utf-8"))
    base = deck_path.parent
    presentation = Presentation()
    presentation.slide_width = Inches(W)
    presentation.slide_height = Inches(H)
    blank = presentation.slide_layouts[6]
    fonts = {"font": deck.get("font", "Gaegu"), "body_font": deck.get("body_font", "Noto Sans KR"), "rough": deck.get("line_mode") == "rough", "Inches": Inches, "Pt": Pt, "RGBColor": RGBColor, "PP_ALIGN": PP_ALIGN, "MSO_SHAPE": MSO_SHAPE}
    layouts = {"title": _layout_title, "content": _layout_content, "full": _layout_full, "quote": _layout_quote, "comparison": _layout_comparison, "timeline": _layout_timeline, "flow": _layout_flow, "table": _layout_table}
    warnings: list[str] = []

    for number, item in enumerate(deck["slides"], 1):
        slide = presentation.slides.add_slide(blank)
        image = None
        if item.get("image"):
            candidate = (base / item["image"]).resolve()
            if candidate.is_file():
                image = _prepared(candidate, max_width, base / ".pptx_cache")
            else:
                warnings.append(f"slide {number}: image not found: {item['image']}")
        if item.get("role"):
            tag = _box(slide, M, 0.4, 6, 0.3, Inches)
            _write(tag, item["role"], font=fonts["font"], size=13, color=ACCENT, Pt=Pt, RGBColor=RGBColor, bold=True)
        layout = item.get("layout", "content")
        if layout not in layouts:
            warnings.append(f"slide {number}: {layout!r} rendered as content")
            layout = "content"
        layouts[layout](slide, item, image, fonts)
        _source(slide, item, fonts["body_font"], Inches, Pt, RGBColor)
        _notes(slide, item)

    output.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output))
    print(f"built: {output} ({len(deck['slides'])} slides, 16:9)")
    for warning in warnings:
        print(f"warning: {warning}")
    return output


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("deck", type=Path)
    parser.add_argument("-o", "--output", type=Path)
    parser.add_argument("--max-image-width", type=int, default=1800)
    args = parser.parse_args()
    deck = args.deck.expanduser().resolve()
    if not deck.is_file():
        print(f"deck not found: {deck}", file=sys.stderr)
        return 2
    try:
        from validate_outline import validate
    except ImportError:
        try:
            from .validate_outline import validate
        except ImportError:
            validate = None
    if validate:
        errors = validate(deck)
        if errors:
            print("outline validation failed:", file=sys.stderr)
            print("\n".join(f"- {error}" for error in errors), file=sys.stderr)
            return 2
    target = (args.output or deck.with_suffix(".pptx")).expanduser().resolve()
    try:
        render(deck, target, args.max_image_width)
    except RuntimeError as exc:
        print(str(exc), file=sys.stderr)
        return 3
    return 0


if __name__ == "__main__":
    sys.exit(main())
